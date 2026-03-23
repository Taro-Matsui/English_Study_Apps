"""
Pick サービス紹介 PowerPoint 生成スクリプト
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

# ─── カラーパレット ───────────────────────────────────────────
AMBER_DARK   = RGBColor(0x92, 0x40, 0x0E)   # amber-800
AMBER_MID    = RGBColor(0xB4, 0x5A, 0x09)   # amber-700
AMBER_LIGHT  = RGBColor(0xFD, 0xF6, 0xEC)   # amber-50
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_900     = RGBColor(0x11, 0x18, 0x27)
GRAY_600     = RGBColor(0x4B, 0x55, 0x63)
EMERALD      = RGBColor(0x05, 0x96, 0x69)   # emerald-600
SLATE_800    = RGBColor(0x1E, 0x29, 0x3B)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank


# ─── ユーティリティ ───────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_w=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    fill = shape.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = shape.line
    if line_rgb:
        line.color.rgb = line_rgb
        line.width = Pt(line_w or 1)
    else:
        line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=GRAY_900,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=14, bold=False, color=GRAY_900,
             align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def slide_bg(slide, rgb=AMBER_LIGHT):
    """スライド背景を塗りつぶす"""
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=rgb)


def header_bar(slide, title, subtitle=None, dark=True):
    """上部ヘッダーバー"""
    fill = AMBER_DARK if dark else AMBER_MID
    add_rect(slide, 0, 0, 13.33, 1.4, fill_rgb=fill)
    add_text(slide, title, 0.4, 0.15, 12, 0.7,
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.82, 12, 0.4,
                 font_size=13, bold=False, color=RGBColor(0xFF, 0xE0, 0xAA),
                 align=PP_ALIGN.LEFT)


def bullet_box(slide, items, left, top, width, height,
               bg=WHITE, title=None, title_color=AMBER_DARK,
               item_size=13.5, item_color=GRAY_900, indent="● "):
    """箇条書きカード"""
    add_rect(slide, left, top, width, height, fill_rgb=bg,
             line_rgb=RGBColor(0xE9, 0xD3, 0xB0), line_w=0.75)
    y_off = top + 0.15
    if title:
        add_text(slide, title, left + 0.18, y_off, width - 0.3, 0.38,
                 font_size=14, bold=True, color=title_color)
        y_off += 0.38
    txBox = slide.shapes.add_textbox(
        Inches(left + 0.18), Inches(y_off),
        Inches(width - 0.3), Inches(height - (y_off - top) - 0.1)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = indent + item
        run.font.size = Pt(item_size)
        run.font.color.rgb = item_color


# ══════════════════════════════════════════════════════════════
# スライド 1 ── タイトル
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=AMBER_DARK)  # 背景
# デコ矩形
add_rect(s, 0, 5.8, 13.33, 1.7, fill_rgb=RGBColor(0x7C, 0x2D, 0x12))
add_rect(s, 0, 0, 0.6, 7.5, fill_rgb=RGBColor(0x7C, 0x2D, 0x12))

add_text(s, "🎸", 0.9, 1.3, 1.5, 1.5, font_size=64, color=WHITE)
add_text(s, "Pick", 2.3, 1.35, 6, 1.2,
         font_size=72, bold=True, color=WHITE)
add_text(s, "ピック", 2.35, 2.5, 4, 0.7, font_size=22, color=RGBColor(0xFF,0xE0,0xAA))
add_text(s, "会話からフレーズをピックして学ぼう",
         0.9, 3.4, 11, 0.8, font_size=26, color=WHITE)
add_text(s, "Pick the words from your real conversations.",
         0.9, 4.0, 11, 0.6, font_size=16, color=RGBColor(0xFF,0xE0,0xAA))
add_text(s, "サービスコンセプト・システム仕様・顧客価値",
         0.9, 5.0, 11, 0.6, font_size=15,
         color=RGBColor(0xFF,0xE0,0xAA))
add_text(s, "2026.03  個人開発 / Railway デプロイ",
         0.9, 6.5, 8, 0.5, font_size=13, color=RGBColor(0xFF,0xC0,0x80))


# ══════════════════════════════════════════════════════════════
# スライド 2 ── 課題とコンセプト
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "Pick が解く課題", "なぜ既存の英語学習アプリでは身につかないのか")

# 左：課題
add_rect(s, 0.3, 1.55, 5.9, 2.9, fill_rgb=RGBColor(0xFF,0xED,0xED),
         line_rgb=RGBColor(0xFC,0xA5,0xA5), line_w=0.75)
add_text(s, "😔  現状の課題", 0.5, 1.65, 5.5, 0.5,
         font_size=15, bold=True, color=RGBColor(0x99,0x1B,0x1B))
problems = [
    "教材フレーズは覚えても実務で使えない",
    "自分の仕事・趣味文脈と乖離している",
    "アプリを続けても「会議で使えた」実感がない",
    "インプット（読む/聴く）と学習が分断されている",
]
bullet_box(s, problems, 0.3, 2.1, 5.9, 2.25,
           bg=RGBColor(0xFF,0xED,0xED),
           title=None, item_size=13, item_color=RGBColor(0x7F,0x1D,0x1D),
           indent="✕ ")

# 右：Pick のアンサー
add_rect(s, 6.6, 1.55, 6.43, 2.9, fill_rgb=RGBColor(0xEC,0xFD,0xF5),
         line_rgb=RGBColor(0x6E,0xE7,0xB7), line_w=0.75)
add_text(s, "✅  Pick のアプローチ", 6.8, 1.65, 6.0, 0.5,
         font_size=15, bold=True, color=RGBColor(0x06,0x54,0x3A))
picks = [
    "自分の議事録・字幕・技術記事が教材になる",
    "AI が「あなたの目的」に合うフレーズを自動抽出",
    "チャレンジで即日確認 → 記憶に定着",
    "ピック〜チャレンジが1アプリで完結",
]
bullet_box(s, picks, 6.6, 2.1, 6.43, 2.25,
           bg=RGBColor(0xEC,0xFD,0xF5),
           title=None, item_size=13, item_color=RGBColor(0x06,0x54,0x3A),
           indent="✓ ")

# 矢印ラベル
add_text(s, "→", 6.0, 2.65, 0.7, 0.7, font_size=36, bold=True,
         color=AMBER_DARK, align=PP_ALIGN.CENTER)

# コンセプト一文
add_rect(s, 0.3, 4.65, 12.73, 1.1, fill_rgb=AMBER_DARK)
add_text(s,
         "「 出会った英語をそのまま学習素材にする 」— 体験と学習の距離をゼロにする",
         0.5, 4.78, 12.3, 0.8,
         font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# フロー図
flow_items = ["📋 テキスト貼付", "🤖 AI フレーズ抽出", "🎯 チャレンジ", "📈 定着・記録"]
positions = [0.35, 3.45, 6.55, 9.65]
for i, (label, x) in enumerate(zip(flow_items, positions)):
    add_rect(s, x, 5.95, 2.9, 0.9, fill_rgb=WHITE,
             line_rgb=AMBER_MID, line_w=1)
    add_text(s, label, x, 5.98, 2.9, 0.85,
             font_size=14.5, bold=(i == 0 or i == 2),
             color=AMBER_DARK, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text(s, "→", x + 2.9, 6.15, 0.55, 0.55,
                 font_size=22, bold=True, color=AMBER_MID, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# スライド 3 ── 顧客価値
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "顧客価値（Customer Value）", "誰に・何を・どうやって届けるか")

segments = [
    {
        "title": "🧑‍💻 エンジニア",
        "who":   "海外チームとの連携・技術文書読解が必要",
        "value": "会議録・PR コメントから即ピック",
        "kpi":   "「あの表現また出た」を学習機会に変換",
        "color_bg": RGBColor(0xEE,0xF2,0xFF),
        "color_bd": RGBColor(0xA5,0xB4,0xFC),
        "color_tt": RGBColor(0x31,0x30,0x8A),
    },
    {
        "title": "💼 ビジネスパーソン",
        "who":   "会議・プレゼン・メール英語を高めたい",
        "value": "議事録・字幕をそのままソースに追加",
        "kpi":   "業務直結フレーズを翌日から使える",
        "color_bg": RGBColor(0xF0,0xFD,0xFA),
        "color_bd": RGBColor(0x99,0xF6,0xE4),
        "color_tt": RGBColor(0x13,0x4E,0x4A),
    },
    {
        "title": "📚 英語学習者（趣味）",
        "who":   "洋画・読書・ポッドキャスト好き",
        "value": "字幕・記事からフレーズをピック",
        "kpi":   "好きなコンテンツが最強の教材になる",
        "color_bg": RGBColor(0xFF,0xF7,0xED),
        "color_bd": RGBColor(0xFD,0xBA,0x74),
        "color_tt": RGBColor(0x7C,0x2D,0x12),
    },
]

xs = [0.3, 4.55, 8.8]
for seg, x in zip(segments, xs):
    add_rect(s, x, 1.55, 4.13, 4.2, fill_rgb=seg["color_bg"],
             line_rgb=seg["color_bd"], line_w=1)
    add_text(s, seg["title"], x + 0.15, 1.65, 3.8, 0.55,
             font_size=16, bold=True, color=seg["color_tt"])
    rows = [
        ("対象ユーザー", seg["who"]),
        ("提供価値",     seg["value"]),
        ("KPI イメージ", seg["kpi"]),
    ]
    y = 2.25
    for label, val in rows:
        add_text(s, label, x + 0.15, y, 3.8, 0.3,
                 font_size=11, bold=True, color=seg["color_tt"])
        add_text(s, val, x + 0.15, y + 0.28, 3.8, 0.55,
                 font_size=13, color=GRAY_900)
        y += 0.9

# 共通価値
add_rect(s, 0.3, 5.9, 12.73, 1.0, fill_rgb=AMBER_DARK)
add_text(s,
         "共通の価値：「自分ごと化」した素材 × AI 採点 × 連続日数（Pick Streak）で学習継続をサポート",
         0.5, 6.0, 12.3, 0.75,
         font_size=14.5, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# スライド 4 ── システム全体像
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "システム全体像", "技術スタックとデータフロー")

layers = [
    ("フロントエンド",  "Next.js 14 App Router / TypeScript / Tailwind CSS / shadcn/ui",        RGBColor(0xEE,0xF2,0xFF), RGBColor(0x81,0x8C,0xF8)),
    ("バックエンド",    "Next.js Route Handler / Claude API（fetch直接）/ Supabase Auth + RLS", RGBColor(0xEC,0xFD,0xF5), RGBColor(0x34,0xD3,0x99)),
    ("データストア",    "Supabase PostgreSQL（RLS マルチユーザー）/ Storage（share-images）",   RGBColor(0xFF,0xF7,0xED), RGBColor(0xFB,0xBF,0x24)),
    ("インフラ",        "Railway（master push → 自動 CI/CD）/ Edge Runtime（OG 画像生成）",     RGBColor(0xF5,0xF3,0xFF), RGBColor(0xA7,0x8B,0xFA)),
    ("AI モデル",       "Sonnet 4.6（フレーズ抽出 8192tok）/ Haiku 4.5（採点 300tok・解説 600tok）", RGBColor(0xFD,0xF2,0xF8), RGBColor(0xF4,0x72,0xB6)),
]

y = 1.6
for layer_name, detail, bg, bd in layers:
    add_rect(s, 0.3, y, 12.73, 0.88, fill_rgb=bg, line_rgb=bd, line_w=0.75)
    add_text(s, layer_name, 0.45, y + 0.07, 2.2, 0.72,
             font_size=13, bold=True, color=GRAY_900)
    add_rect(s, 2.6, y + 0.12, 0.03, 0.6, fill_rgb=bd)
    add_text(s, detail, 2.75, y + 0.12, 10.2, 0.65,
             font_size=12.5, color=GRAY_600)
    y += 0.97

# セキュリティ補足
add_rect(s, 0.3, 6.5, 12.73, 0.7, fill_rgb=RGBColor(0xFF,0xF1,0xF2),
         line_rgb=RGBColor(0xFC,0xA5,0xA5), line_w=0.5)
add_text(s,
         "🔒 セキュリティ：UUID バリデーション / 入力長制限 / allowlist / PostgREST インジェクション対策 / CSP / HSTS",
         0.5, 6.55, 12.3, 0.6, font_size=12, color=RGBColor(0x99,0x1B,0x1B))


# ══════════════════════════════════════════════════════════════
# スライド 5 ── 主要機能
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "主要機能一覧", "実装済みフィーチャー（2026-03 時点）")

features = [
    ("📋", "ソース追加\n（ピック）",     "テキスト貼付 → Claude Sonnet が学習目的別にフレーズ自動抽出（JSON 部分回復付き）"),
    ("🎯", "チャレンジ",               "テキスト入力 → Claude Haiku でAI採点。正誤・惜しいの3段階判定"),
    ("💡", "AI Coach 解説",           "不正解・惜しい時に自動表示。話す/聴く/使い分け/NGパターンの4セクション"),
    ("📊", "Progress & Streak",       "14日チャート・難易度別精度・連続日数カレンダー（過去6ヶ月）"),
    ("🔄", "ピックアップ チャレンジ",   "直近5セッションで2回以上不正解のフレーズを優先出題（focus モード）"),
    ("𝕏",  "X シェア",                "Canvas 1200×630 PNG 生成 → Supabase Storage → Twitter Cards 表示（Safari対応）"),
    ("📱", "PWA",                     "iOS Safari / Android Chrome インストール対応。ホーム画面からオフライン起動"),
    ("🌐", "多言語 UI",               "JA/EN 切り替え（i18n.tsx）。学習内容は英語固定"),
    ("🎓", "オンボーディング",          "目的・レベル・英語力チェック（多肢選択）→ シードフレーズ自動挿入"),
]

cols = 3
rows_count = (len(features) + cols - 1) // cols
xs_f = [0.3, 4.55, 8.8]
ys_f = [1.6, 3.15, 4.7]

for idx, (icon, title, desc) in enumerate(features):
    col = idx % cols
    row = idx // cols
    x = xs_f[col]
    y = ys_f[row]
    add_rect(s, x, y, 4.13, 1.35, fill_rgb=WHITE,
             line_rgb=RGBColor(0xE9,0xD3,0xB0), line_w=0.75)
    add_text(s, icon + "  " + title, x + 0.15, y + 0.08, 3.8, 0.45,
             font_size=13.5, bold=True, color=AMBER_DARK)
    add_text(s, desc, x + 0.15, y + 0.52, 3.8, 0.75,
             font_size=11.5, color=GRAY_600)

add_text(s, "※ 課金プラン (Pick Free/Starter/Pro) は設計済み・未実装",
         0.3, 6.2, 12.73, 0.4, font_size=11.5, color=GRAY_600)


# ══════════════════════════════════════════════════════════════
# スライド 6 ── データフロー詳細（ピック → チャレンジ）
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "データフロー詳細", "ピック（フレーズ抽出）→ チャレンジ（AI採点）の流れ")

steps_pick = [
    ("1", "テキスト入力",    "User"),
    ("2", "Claude Sonnet 抽出", "AI"),
    ("3", "フレーズ保存",    "DB"),
    ("4", "マイピックリスト表示", "UI"),
]
steps_quiz = [
    ("A", "フレーズ出題",    "DB"),
    ("B", "テキスト回答",    "User"),
    ("C", "Claude Haiku 採点", "AI"),
    ("D", "解説 + 記録保存",  "DB"),
]

def draw_flow(slide, steps, top, label, label_color):
    add_text(slide, label, 0.3, top - 0.38, 2.5, 0.35,
             font_size=13, bold=True, color=label_color)
    col_w = 12.73 / len(steps)
    for i, (num, title, tag) in enumerate(steps):
        x = 0.3 + i * col_w
        add_rect(slide, x + 0.05, top, col_w - 0.15, 1.0,
                 fill_rgb=WHITE, line_rgb=AMBER_MID, line_w=0.75)
        add_text(slide, num, x + 0.1, top + 0.03, 0.45, 0.38,
                 font_size=18, bold=True, color=AMBER_DARK)
        add_text(slide, title, x + 0.5, top + 0.05, col_w - 0.65, 0.45,
                 font_size=13, bold=True, color=GRAY_900)
        tag_colors = {"User": EMERALD, "AI": RGBColor(0x7C,0x3A,0xED), "DB": AMBER_MID}
        add_rect(slide, x + 0.5, top + 0.58, 0.8, 0.28,
                 fill_rgb=tag_colors.get(tag, GRAY_600))
        add_text(slide, tag, x + 0.5, top + 0.57, 0.8, 0.3,
                 font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_text(slide, "→", x + col_w - 0.12, top + 0.3, 0.4, 0.45,
                     font_size=20, bold=True, color=AMBER_MID)

draw_flow(s, steps_pick, 1.85, "📋 ピックフロー", AMBER_DARK)
draw_flow(s, steps_quiz, 3.55, "🎯 チャレンジフロー", EMERALD)

# 補足メモ
notes = [
    "・JSON 部分回復: max_tokens 超過でも 3 段階フォールバックで切断 JSON を復元",
    "・新規フレーズ優先: 直近7日以内 → 通常 → quickMastered（直近10セッションで3回以上 <3.5sec 正解）の順で出題",
    "・response_time_ms 計測: 回答速度を quiz_answers テーブルに保存し quickMastered 判定に使用",
    "・論理削除: phrases.deleted_at。全クエリで .is('deleted_at', null) フィルタ必須",
]
y_note = 4.95
add_rect(s, 0.3, y_note, 12.73, 2.2, fill_rgb=RGBColor(0xF8,0xF5,0xF0),
         line_rgb=RGBColor(0xE9,0xD3,0xB0), line_w=0.5)
add_text(s, "📝 実装メモ", 0.5, y_note + 0.08, 3, 0.35,
         font_size=12, bold=True, color=AMBER_DARK)
for i, note in enumerate(notes):
    add_text(s, note, 0.5, y_note + 0.42 + i * 0.42, 12.4, 0.4,
             font_size=11.5, color=GRAY_600)


# ══════════════════════════════════════════════════════════════
# スライド 7 ── 認証・セキュリティ・DB
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "認証・DB・セキュリティ設計", "マルチユーザー対応の基盤")

# 認証フロー
add_rect(s, 0.3, 1.55, 6.1, 3.3, fill_rgb=WHITE,
         line_rgb=RGBColor(0xE9,0xD3,0xB0), line_w=0.75)
add_text(s, "🔐 認証フロー", 0.5, 1.65, 5.5, 0.4,
         font_size=14, bold=True, color=AMBER_DARK)
auth_steps = [
    "/login → signUp()  →  確認メール送信",
    "→ /auth/callback → exchangeCodeForSession()",
    "→ middleware: onboarding_complete チェック",
    "→ 未完了 → /onboarding（目的・レベル・英語力）",
    "→ refreshSession() 必須（JWT 更新）→ router.push(/)",
    "→ WelcomeGuide（初回）→ HintBubble → 通常利用",
]
y = 2.1
for step in auth_steps:
    add_text(s, step, 0.5, y, 5.7, 0.35, font_size=11.5, color=GRAY_900)
    y += 0.35

# DB スキーマ
add_rect(s, 6.6, 1.55, 6.43, 3.3, fill_rgb=WHITE,
         line_rgb=RGBColor(0xE9,0xD3,0xB0), line_w=0.75)
add_text(s, "🗄️ DB スキーマ（主要テーブル）", 6.78, 1.65, 6.0, 0.4,
         font_size=14, bold=True, color=AMBER_DARK)
db_items = [
    "phrases: id / user_id / phrase / translation",
    "         explanation / deleted_at / added_date",
    "quiz_sessions: id / user_id / started_at",
    "quiz_answers: session_id / phrase_id / status",
    "              response_time_ms (採点速度記録)",
    "RLS: 全テーブルで user_id = auth.uid() 強制",
    "Admin key + .eq('user_id') でバイパス管理",
]
y = 2.1
for item in db_items:
    add_text(s, item, 6.78, y, 6.1, 0.38, font_size=11.5, color=GRAY_900)
    y += 0.38

# セキュリティ
sec_items = [
    "UUID v4 バリデーション（phrases/[id]・jobs/[id]）",
    "入力長制限: phrase(200) / answer(500) / domain(100) / file(2MB, 200k chars)",
    "allowlist: source_type / delete_reason / study_purpose / study_level",
    "PostgREST インジェクション対策: .or() 文字列から特殊文字をストリップ",
    "SSRF 対策: import URL のプライベート IP 拒否（isAllowedHost）",
    "CSP: connect/img-src に Supabase・Google・GTM を明示。Storage 画像は *.supabase.co 必須",
]
add_rect(s, 0.3, 4.95, 12.73, 2.25, fill_rgb=RGBColor(0xFF,0xF1,0xF2),
         line_rgb=RGBColor(0xFC,0xA5,0xA5), line_w=0.75)
add_text(s, "🔒 セキュリティ実装", 0.5, 5.05, 5, 0.4,
         font_size=14, bold=True, color=RGBColor(0x99,0x1B,0x1B))
y = 5.48
for item in sec_items:
    add_text(s, "  ・" + item, 0.5, y, 12.3, 0.32, font_size=11.5,
             color=RGBColor(0x7F,0x1D,0x1D))
    y += 0.32


# ══════════════════════════════════════════════════════════════
# スライド 8 ── X シェア & PWA
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "X シェア & PWA", "Safari 対策・Twitter Cards・インストール体験")

# X share flow
add_rect(s, 0.3, 1.55, 12.73, 2.0, fill_rgb=WHITE,
         line_rgb=RGBColor(0xE9,0xD3,0xB0), line_w=0.75)
add_text(s, "𝕏 シェアフロー（Safari 対策込み）", 0.5, 1.65, 8, 0.4,
         font_size=14, bold=True, color=AMBER_DARK)
xsteps = [
    "① チャレンジ完了",
    "② 即・Canvas PNG 生成\n& Storage アップロード",
    "③ シェアボタン押下\nwindow.open()（同期）",
    "④ Twitter が /share/{id}\nをクロール",
    "⑤ og:image → Storage URL\nCard 表示",
]
for i, step in enumerate(xsteps):
    x = 0.4 + i * 2.53
    add_rect(s, x, 2.08, 2.35, 1.25, fill_rgb=AMBER_LIGHT,
             line_rgb=AMBER_MID, line_w=0.5)
    add_text(s, step, x + 0.1, 2.13, 2.15, 1.15,
             font_size=11.5, color=GRAY_900)
    if i < 4:
        add_text(s, "→", x + 2.35, 2.5, 0.25, 0.45,
                 font_size=16, bold=True, color=AMBER_MID)

# Safari ポイント
add_rect(s, 0.3, 3.7, 6.1, 1.85, fill_rgb=RGBColor(0xFF,0xF1,0xF2),
         line_rgb=RGBColor(0xFC,0xA5,0xA5), line_w=0.75)
add_text(s, "⚠️ Safari 特有の制約", 0.5, 3.8, 5.5, 0.4,
         font_size=13, bold=True, color=RGBColor(0x99,0x1B,0x1B))
safari = [
    "await を挟むと window.open() がブロックされる",
    "window.open() 後はタブが suspend → バックグラウンド処理 kill",
    "解決：チャレンジ完了時（FG）に先行アップロード完了",
]
y = 4.24
for s_item in safari:
    add_text(s, "✕ " + s_item, 0.5, y, 5.7, 0.38,
             font_size=11.5, color=RGBColor(0x7F,0x1D,0x1D))
    y += 0.38

# PWA
add_rect(s, 6.6, 3.7, 6.43, 1.85, fill_rgb=RGBColor(0xEC,0xFD,0xF5),
         line_rgb=RGBColor(0x34,0xD3,0x99), line_w=0.75)
add_text(s, "📱 PWA 対応", 6.78, 3.8, 5.8, 0.4,
         font_size=13, bold=True, color=RGBColor(0x06,0x54,0x3A))
pwa_items = [
    "app/icon.png（静的ファイル）→ PWA + favicon 兼用",
    "iOS Safari: 「共有 → ホーム画面に追加」バナー案内",
    "Android Chrome: beforeinstallprompt でネイティブ誘導",
    "manifest.ts: name = 'Pick' / short_name = 'Pick'",
]
y = 4.24
for p_item in pwa_items:
    add_text(s, "✓ " + p_item, 6.78, y, 6.1, 0.38,
             font_size=11.5, color=RGBColor(0x06,0x54,0x3A))
    y += 0.38

# Twitter Cards
add_rect(s, 0.3, 5.65, 12.73, 1.1, fill_rgb=AMBER_LIGHT,
         line_rgb=AMBER_MID, line_w=0.5)
add_text(s, "🐦 Twitter Cards 要件: summary_large_image / 1200×630px 横長必須 / "
            "/share/{id} & /api/og は PUBLIC_PATHS に含める（Twitter クローラーは未認証）",
         0.5, 5.75, 12.3, 0.85, font_size=12.5, color=GRAY_900)


# ══════════════════════════════════════════════════════════════
# スライド 9 ── ブランド & UX
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "ブランド & UX デザイン", "Pick Language Design v1.0")

# カラーテーマ
add_rect(s, 0.3, 1.55, 5.0, 2.5, fill_rgb=WHITE,
         line_rgb=RGBColor(0xE9,0xD3,0xB0), line_w=0.75)
add_text(s, "🎨 ウォームベージュテーマ", 0.5, 1.65, 4.5, 0.4,
         font_size=13, bold=True, color=AMBER_DARK)
colors_disp = [
    ("bg-amber-50",  RGBColor(0xFD,0xF6,0xEC), "#FDF6EC  ベース背景"),
    ("bg-white/90",  RGBColor(0xFF,0xFF,0xFF),  "#FFFFFF  カード"),
    ("bg-amber-800", RGBColor(0x92,0x40,0x0E),  "#92400E  ボタン Primary"),
    ("bg-emerald-600", EMERALD,                 "#059669  ボタン CTA"),
]
y = 2.1
for cls, rgb, label in colors_disp:
    add_rect(s, 0.5, y, 0.45, 0.32, fill_rgb=rgb,
             line_rgb=RGBColor(0xCC,0xCC,0xCC), line_w=0.5)
    add_text(s, label, 1.05, y, 4.0, 0.32, font_size=11.5, color=GRAY_900)
    y += 0.4
add_text(s, "dark: クラス全廃。ThemeProvider は no-op。",
         0.5, 3.65, 4.5, 0.3, font_size=10.5, color=GRAY_600)

# 用語
add_rect(s, 5.6, 1.55, 7.43, 2.5, fill_rgb=WHITE,
         line_rgb=RGBColor(0xE9,0xD3,0xB0), line_w=0.75)
add_text(s, "📖 ブランド用語（抜粋）", 5.78, 1.65, 7.0, 0.4,
         font_size=13, bold=True, color=AMBER_DARK)
terms = [
    ("ピックする",             "フレーズを選び取る行為（動詞）"),
    ("チャレンジ",             "クイズ画面 /quiz"),
    ("マイピックリスト",        "ピック済みフレーズ一覧 /phrases"),
    ("ピックアップ チャレンジ", "苦手フレーズ復習モード（focus）"),
    ("連続日数",               "学習継続日数 /streak"),
    ("Add a Source",          "テキスト追加 /library/import"),
]
y = 2.1
for term, desc in terms:
    add_text(s, term, 5.78, y, 3.0, 0.35, font_size=12, bold=True, color=AMBER_DARK)
    add_text(s, desc, 8.7, y, 4.1, 0.35, font_size=11.5, color=GRAY_600)
    y += 0.38

# モバイル UX
add_rect(s, 0.3, 4.2, 12.73, 2.05, fill_rgb=WHITE,
         line_rgb=RGBColor(0xE9,0xD3,0xB0), line_w=0.75)
add_text(s, "📱 モバイル UX ガイドライン", 0.5, 4.3, 8, 0.4,
         font_size=13, bold=True, color=AMBER_DARK)
ux_items = [
    "input の font-size は必ず 16px（iOS Safari 自動ズーム防止）",
    "min-h-screen (100vh) を使用。min-h-[100dvh] はキーボード展開時にレイアウトシフト",
    "sticky ヘッダー + 通常フロー。flex split layout は NG（iOS keyboard 展開で問題発生）",
    "step 変化時: question → scrollTo(top:0, instant) + setTimeout focus(80ms)",
    "戻るボタン: text-2xl p-2 -ml-2（約40px タップ領域確保）",
]
y = 4.72
for ux in ux_items:
    add_text(s, "● " + ux, 0.5, y, 12.3, 0.32, font_size=11.5, color=GRAY_900)
    y += 0.32


# ══════════════════════════════════════════════════════════════
# スライド 10 ── 現状と今後
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "現状と今後のロードマップ", "フェーズ・KPI・次のステップ")

# 現状
add_rect(s, 0.3, 1.55, 5.9, 3.2, fill_rgb=WHITE,
         line_rgb=RGBColor(0xE9,0xD3,0xB0), line_w=0.75)
add_text(s, "📍 現状（2026-03）", 0.5, 1.65, 5.5, 0.4,
         font_size=14, bold=True, color=AMBER_DARK)
status_items = [
    "テストユーザーのみで運用中",
    "本番 URL: usepick.win",
    "AdSense 審査保留中（privacy/terms/about/contact 整備済み）",
    "マイグレーション 001〜013 適用済み",
    "課金プラン設計済み・未実装",
    "本番ユーザー獲得前フェーズ",
]
y = 2.1
for item in status_items:
    add_text(s, "・" + item, 0.5, y, 5.6, 0.38, font_size=12.5, color=GRAY_900)
    y += 0.38

# 次のステップ
add_rect(s, 6.6, 1.55, 6.43, 3.2, fill_rgb=WHITE,
         line_rgb=RGBColor(0xE9,0xD3,0xB0), line_w=0.75)
add_text(s, "🚀 次のステップ", 6.78, 1.65, 6.0, 0.4,
         font_size=14, bold=True, color=EMERALD)
next_items = [
    "AdSense 審査通過 → 広告収益化",
    "課金プラン実装（Pick Starter/Pro）",
    "X / SNS での認知拡大",
    "オーディオ再生 TTS 対応",
    "フレーズ CSV エクスポート",
    "チーム共有機能（複数人でソース共有）",
]
y = 2.1
for item in next_items:
    add_text(s, "→ " + item, 6.78, y, 6.1, 0.38, font_size=12.5,
             color=EMERALD)
    y += 0.38

# 差別化
add_rect(s, 0.3, 4.88, 12.73, 2.28, fill_rgb=AMBER_DARK)
add_text(s, "💎  Pick の差別化ポイント", 0.5, 4.98, 12.3, 0.5,
         font_size=15, bold=True, color=WHITE)
diffs = [
    "① 「自分の体験テキスト」が素材 — 汎用教材アプリとの根本差別化",
    "② AI 採点 × 解説 × Streak — インプット・アウトプット・継続が1アプリで完結",
    "③ モバイルファースト PWA — iOS Safari まで完全対応",
    "④ 個人開発でありながら エンタープライズ水準のセキュリティ（RLS・CSP・SSRF対策）",
]
y = 5.5
for diff in diffs:
    add_text(s, diff, 0.5, y, 12.3, 0.3, font_size=12, color=WHITE)
    y += 0.32


# ─── 保存 ─────────────────────────────────────────────────────
out = r"c:\Users\04595\engineer-english-app\docs\Pick_Service_Overview.pptx"
prs.save(out)
print(f"Saved: {out}")
